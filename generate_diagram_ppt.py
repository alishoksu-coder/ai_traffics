import os
import zlib
import base64
import requests
from pptx import Presentation
from pptx.util import Inches

def encode_kroki(text):
    compressed = zlib.compress(text.encode('utf-8'), 9)
    return base64.urlsafe_b64encode(compressed).decode('utf-8')

diagrams = {
    "Architecture": """flowchart LR
    classDef default fill:#e3e8ec,stroke:#b0bec5,stroke-width:1px,color:#37474f,rx:6px,ry:6px;
    classDef layer fill:#cfd8dc,stroke:#90a4ae,stroke-width:2px,color:#263238,rx:8px,ry:8px;
    classDef diamond fill:#cfd8dc,stroke:#b0bec5,stroke-width:1px,color:#37474f;
    
    Start([Start: Notes & Mood Tracking App]) --> Layer{Which Layer?}:::diamond
    
    Layer --> P[Presentation]:::layer
    Layer --> D[Domain]:::layer
    Layer --> Da[Data]:::layer
    
    P --> S1{Select Screen}:::diamond
    S1 --> |View Notes| NG[Notes Grid Screen]
    S1 --> |Edit Notes| NE[Note Editor]
    S1 --> |Calendar View| CS[Calendar Screen]
    S1 --> |Mood Tracking| MD[Mood Diary]
    S1 --> |Analytics| NGr[Note Graph]
    
    D --> S2{Select Service}:::diamond
    S2 --> |AI Processing| GS[Gemini Service AI Analysis]
    S2 --> |Mood Analysis| MC[Mood Color NLP Processing]
    S2 --> |Task Extraction| TP[Task Planner NLP Parser]
    S2 --> |Journaling| JS[Journal Streak Tracking]
    S2 --> |Music Therapy| MM[Mood Music Service]
    
    Da --> S3{Select Storage}:::diamond
    S3 --> |Notes Storage| NR[Notes Repository Hive NoSQL]
    S3 --> |Chat History| CH[AI Chat Storage Hive]
    S3 --> |Settings| PS[Pin Service SharedPrefs]
    S3 --> |Widgets| HW[Home Widget Service]
    
    NG & NE & CS & MD & NGr & GS & MC & TP & JS & MM & NR & CH & PS & HW --> Display([Display Notes])""",
    
    "Lifecycle": """flowchart TD
    classDef default fill:#b35c5c,stroke:#ffffff,stroke-width:1px,color:#ffffff,rx:4px,ry:4px;
    
    Create[Жасалу<br>Create] -->|transition| Active[active]
    
    Active -->|archive| Archived[archived]
    Active -->|action| Action1[Өңдеу<br>Edit: title, content, tags...]
    Active -->|action| Action2[MoodColor.detect<br>Update color, accentColor]
    Active -->|action| Action3[TaskPlannerParser.parse<br>Update scheduledAt...]
    Active -->|action| Action4[VoiceRecorder<br>Update audioPath]
    
    Action1 & Action2 & Action3 & Action4 -->|back to| Active
    
    Archived -->|delete| Deleted[deleted]
    Archived -->|restore| Restore[Қалпына келтіру<br>Restore]
    Restore --> Active
    
    Deleted -->|permanent| Permanent[Жойылу<br>Delete]""",
    
    "MoodColor": """flowchart TD
    classDef default fill:#1a4cd6,stroke:#ffffff,stroke-width:1px,color:#ffffff,rx:4px,ry:4px;
    classDef check fill:#1a4cd6,stroke:#ffffff,stroke-width:1px,color:#ffffff;
    
    Start[Кіріс мәтін] --> Check{tokens.isEmpty?}:::check
    
    Check -->|Иә| Neutral[Mood.neutral]
    Check -->|Жоқ| Token[Токенизация RegEx]
    
    Token --> Scan[Кілт сөздер сканерлеу]
    
    Scan --> Pos[pos x2]
    Scan --> Urg[urg x3]
    Scan --> Sad[sad x2]
    
    Pos & Urg & Sad --> Max[Максималды ұпай]
    
    Max --> Best[bestMood анықтау]
    Best --> Colors[Primary + Accent түстер]
    Colors --> Soft[SoftGradients генерация]
    Soft --> Final[Карточка градиенті]""",
    
    "NoteEditorScreen": """flowchart TD
    classDef default fill:#f74ce5,stroke:#ffffff,stroke-width:1px,color:#ffffff,rx:4px,ry:4px;
    classDef action fill:#f74ce5,stroke:#ffffff,stroke-width:1px,color:#ffffff;
    
    Start[NoteEditorScreen] --> Action{What action?}:::action
    
    Action -->|Detect Mood| M1[MoodColor.detect]
    M1 --> M2[Apply SoftGradients карточка түсі]
    
    Action -->|Parse Tasks| T1[TaskPlannerParser.parse]
    T1 --> T2[Note.copyWith schedule, location...]
    
    Action -->|Update Music| Mus1[MoodMusicService.updateMood]
    Mus1 --> Mus2[AudioPlayer x2 кроссфейдинг]
    
    M2 & T2 & Mus2 --> Save{Save to Repository?}:::action
    
    Save -->|Yes| S1[NotesRepo.upsert]
    S1 --> S2[Hive NoSQL]
    S2 --> S3[HomeWidgetService.updateData]
    S3 --> End1[Update Complete]
    
    Save -->|No| C1[Cancel]
    C1 --> C2[Discard Changes]""",
    
    "PIN_Verify": """flowchart TD
    classDef pink fill:#f74ce5,stroke:#ffffff,stroke-width:1px,color:#ffffff,rx:4px,ry:4px;
    classDef check fill:#f74ce5,stroke:#ffffff,stroke-width:1px,color:#ffffff;
    
    Input[Пайдаланушы PIN енгізеді]:::pink --> GetSalt[SharedPreferences-тен salt алу]:::pink
    GetSalt --> Hash[HMAC-SHA256<br>key=saved_salt<br>data=input_pin]:::pink
    Hash --> Check{computed_hash == saved_hash ?}:::check
    
    Check -->|иә| Success[Кіру / Login]:::pink
    Check -->|жоқ| Fail[Қате анимация / Error]:::pink"""
}

# Create a presentation
prs = Presentation()

for name, code in diagrams.items():
    print(f"Generating image for {name}...")
    encoded = encode_kroki(code)
    url = f"https://kroki.io/mermaid/png/{encoded}"
    
    # Download image
    response = requests.get(url)
    if response.status_code == 200:
        filename = f"diagram_{name}.png"
        with open(filename, "wb") as f:
            f.write(response.content)
            
        # Add to presentation
        slide_layout = prs.slide_layouts[5] # blank with title
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = name
        
        # Add picture (center it roughly)
        try:
            # We don't know the exact dimensions, we will fit it within the slide
            slide.shapes.add_picture(filename, Inches(1), Inches(1.5), width=Inches(8))
        except Exception as e:
            print(f"Error adding {filename} to ppt: {e}")
            
    else:
        print(f"Failed to generate {name}. HTTP {response.status_code}")

ppt_filename = "Diagrams_Screenshots.pptx"
prs.save(ppt_filename)
print(f"Saved all diagrams to {ppt_filename}")
