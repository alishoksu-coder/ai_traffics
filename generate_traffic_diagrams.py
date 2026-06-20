import os
import zlib
import base64
import requests
from pptx import Presentation
from pptx.util import Inches

def encode_kroki(text):
    compressed = zlib.compress(text.encode('utf-8'), 9)
    return base64.urlsafe_b64encode(compressed).decode('utf-8')

diagrams = {
    "Функциональная диаграмма": """flowchart TD
    A["🚀 Splash Screen"] --> B["🔐 Аутентификация"]
    B --> B1["FaceID / TouchID"]
    B --> B2["PIN-код"]
    B --> B3["Админ-логин"]
    
    B1 & B2 --> C["🏠 Главный экран"]
    B3 --> ADM["🛡️ Админ-панель"]

    C --> T1["🗺️ Карта"]
    C --> T2["🧭 Навигатор"]
    C --> T3["🚗 Режим вождения"]
    C --> T4["📊 Метрики"]
    C --> T5["⚙️ Ещё"]

    T1 --> T1a["Тепловая карта трафика"]
    T1 --> T1b["144 точки мониторинга"]
    T1 --> T1c["Цветовая индикация загруженности"]
    T1 --> T1d["Кастомные маркеры"]

    T2 --> T2a["🔍 Поиск маршрута A→B"]
    T2 --> T2b["♿ Инклюзивный маршрут"]
    T2 --> T2c["🎤 Голосовой ввод"]
    T2 --> T2d["⏱️ ETA с учётом трафика"]
    T2 --> T2e["🧠 AI-рекомендации"]

    T3 --> T3a["GPS-трекинг в реальном времени"]
    T3 --> T3b["Спидометр"]
    T3 --> T3c["Пошаговые подсказки"]
    T3 --> T3d["Антистресс-режим"]

    T4 --> T4a["LSTM прогноз загруженности"]
    T4 --> T4b["Графики трафика"]
    T4 --> T4c["Z-Score аномалии"]
    T4 --> T4d["Погодное влияние"]

    T5 --> T5a["👥 Друзья"]
    T5 --> T5b["💡 Советы"]
    T5 --> T5c["🔒 Настройки безопасности"]
    T5 --> T5d["🌙 Тёмная тема"]
    T5 --> T5e["📡 Сегменты дорог"]

    ADM --> ADM1["Мониторинг всех сегментов"]
    ADM --> ADM2["Управление инцидентами"]
    ADM --> ADM3["Статистика и отчёты"]

    style A fill:#3b82f6,color:#fff,stroke:none
    style B fill:#7c3aed,color:#fff,stroke:none
    style C fill:#2563eb,color:#fff,stroke:none
    style ADM fill:#ef4444,color:#fff,stroke:none
    style T1 fill:#10b981,color:#fff,stroke:none
    style T2 fill:#f59e0b,color:#fff,stroke:none
    style T3 fill:#06b6d4,color:#fff,stroke:none
    style T4 fill:#8b5cf6,color:#fff,stroke:none
    style T5 fill:#64748b,color:#fff,stroke:none""",
    
    "Архитектура LSTM": """flowchart TD
      A["📡 IoT-сенсоры (144 сегмента)"] --> B["🗄️ База данных (PostgreSQL)"]
      B --> C["🧹 Препроцессинг (MinMaxScaler)"]
      
      subgraph LSTM ["🧠 LSTM Нейронная сеть (PyTorch)"]
         direction TB
         L1["Входные данные<br>(Скорость, Время, Погода, Праздники)"] --> L2["LSTM Cell 1<br>(Память долгих зависимостей)"]
         L2 --> L3["LSTM Cell 2<br>(Выявление скрытых паттернов)"]
         L3 --> L4["Fully Connected Layer<br>(Финальная регрессия)"]
      end
      
      C --> LSTM
      LSTM --> D["🔮 Прогноз трафика<br>(на 30-60 минут вперёд)"]
      D --> E["📱 Вывод в Mobile / Web"]
      
      style A fill:#0ea5e9,color:#fff,stroke:none
      style B fill:#334155,color:#fff,stroke:none
      style C fill:#f59e0b,color:#fff,stroke:none
      style LSTM fill:#1e1b4b,color:#fff,stroke:#8b5cf6,stroke-width:2px
      style L1 fill:#4c1d95,color:#fff,stroke:none
      style L2 fill:#5b21b6,color:#fff,stroke:none
      style L3 fill:#6d28d9,color:#fff,stroke:none
      style L4 fill:#7c3aed,color:#fff,stroke:none
      style D fill:#ec4899,color:#fff,stroke:none
      style E fill:#10b981,color:#fff,stroke:none"""
}

# Create a presentation
try:
    prs = Presentation('Diagrams_Screenshots.pptx')
except:
    prs = Presentation()

for name, code in diagrams.items():
    print(f"Generating image for {name}...")
    encoded = encode_kroki(code)
    url = f"https://kroki.io/mermaid/png/{encoded}"
    
    # Download image
    response = requests.get(url)
    if response.status_code == 200:
        filename = f"diagram_{name}.png"
        with open(filename, "wb") as f:
            f.write(response.content)
            
        # Add to presentation
        slide_layout = prs.slide_layouts[5] # blank with title
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = name
        
        # Add picture (center it roughly)
        try:
            slide.shapes.add_picture(filename, Inches(1), Inches(1.5), width=Inches(8))
        except Exception as e:
            print(f"Error adding {filename} to ppt: {e}")
            
    else:
        print(f"Failed to generate {name}. HTTP {response.status_code}")

ppt_filename = "Diagrams_Screenshots.pptx"
prs.save(ppt_filename)
print(f"Saved AI Traffic diagrams to {ppt_filename}")
