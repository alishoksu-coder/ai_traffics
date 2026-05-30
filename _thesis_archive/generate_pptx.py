from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

def create_premium_pptx():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Colors (Modern Blue Theme matching the website)
    BG_COLOR = RGBColor(241, 245, 249)    # Light slate background
    PRIMARY_TEXT = RGBColor(15, 23, 42)   # Dark slate text
    SECONDARY_TEXT = RGBColor(71, 85, 105) # Muted text
    ACCENT_BLUE = RGBColor(37, 99, 235)   # Vibrant blue
    CARD_BG = RGBColor(255, 255, 255)     # White card
    BORDER_COLOR = RGBColor(226, 232, 240) # Light border

    def apply_premium_style(slide, title_text):
        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = BG_COLOR
        
        # Add a "Glass" Card effect (Large rounded rectangle)
        card = slide.shapes.add_shape(
            1, Inches(0.4), Inches(0.4), Inches(12.5), Inches(6.7)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = BORDER_COLOR
        card.line.width = Pt(1)
        
        # Header Blue bar
        header_bar = slide.shapes.add_shape(
            1, Inches(0.4), Inches(0.4), Inches(12.5), Inches(0.15)
        )
        header_bar.fill.solid()
        header_bar.fill.fore_color.rgb = ACCENT_BLUE
        header_bar.line.fill.background()

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(11), Inches(1))
        tf = title_box.text_frame
        tf.text = title_text
        p = tf.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(32)
        p.font.color.rgb = PRIMARY_TEXT
        
        # Bottom decorative element
        footer_line = slide.shapes.add_shape(
            1, Inches(0.8), Inches(1.6), Inches(1.5), Inches(0.04)
        )
        footer_line.fill.solid()
        footer_line.fill.fore_color.rgb = ACCENT_BLUE
        footer_line.line.fill.background()

        return slide

    def add_title_slide(title_text, subtitle_text):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        # Full screen background with a gradient-like look
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(255, 255, 255)
        
        # Add side accent
        side_bar = slide.shapes.add_shape(1, 0, 0, Inches(0.3), prs.slide_height)
        side_bar.fill.solid()
        side_bar.fill.fore_color.rgb = ACCENT_BLUE
        side_bar.line.fill.background()

        # Main Title
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(2), Inches(11), Inches(2.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.bold = True
        p.font.size = Pt(46)
        p.font.color.rgb = PRIMARY_TEXT
        p.alignment = PP_ALIGN.LEFT

        # Subtitle
        txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(4.8), Inches(11), Inches(1.5))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = subtitle_text
        p2.font.size = Pt(24)
        p2.font.color.rgb = SECONDARY_TEXT
        
        # Add a "Project Site" badge
        badge = slide.shapes.add_shape(1, Inches(0.8), Inches(0.8), Inches(4), Inches(0.5))
        badge.fill.solid()
        badge.fill.fore_color.rgb = RGBColor(239, 246, 255)
        badge.line.color.rgb = RGBColor(191, 219, 254)
        
        b_text = slide.shapes.add_textbox(Inches(1), Inches(0.85), Inches(3.8), Inches(0.4))
        bf = b_text.text_frame
        bp = bf.paragraphs[0]
        bp.text = "AI TRAFFIC PROJECT | 2026"
        bp.font.size = Pt(12)
        bp.font.bold = True
        bp.font.color.rgb = ACCENT_BLUE

    def add_content_slide(title_text, points):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        apply_premium_style(slide, title_text)

        # Body text area
        txBox = slide.shapes.add_textbox(Inches(1), Inches(2.2), Inches(11), Inches(4.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        
        for point in points:
            p = tf.add_paragraph()
            p.text = f"• {point}"
            p.font.size = Pt(22)
            p.font.color.rgb = RGBColor(51, 65, 85)
            p.space_after = Pt(15)
            # Indent
            p.level = 0

    # --- SLIDES CONTENT ---
    add_title_slide(
        "Қалалық ортадағы көлік ағындарын бақылау мен болжауға арналған AI-қосымша әзірлеу",
        "Студент: Сулейменов Алишер\nҒылыми жетекші: Кусаинова Айнур"
    )

    add_content_slide("1. Тақырыптың өзектілігі", [
        "Урбанизацияның өсуі және көлік кептелісі мәселесі.",
        "Қазақстанның ірі қалаларындағы трафик индексінің нашарлауы.",
        "Экологиялық және экономикалық шығындарды азайту қажеттілігі.",
        "Жасанды интеллект көмегімен басқарудың жаңа мүмкіндіктері."
    ])

    add_content_slide("2. Жұмыс мақсаты", [
        "LSTM нейрондық желілері негізінде трафикті болжау жүйесін құру.",
        "Нақты уақыттағы IoT мониторингін іске асыру.",
        "Ыңғайлы мобильді және веб-платформаны әзірлеу.",
        "Трафикті оңтайландыру арқылы жол қауіпсіздігін арттыру."
    ])

    add_content_slide("3. Зерттеу міндеттері", [
        "Қолданыстағы навигациялық жүйелерді талдау (2GIS, Yandex, Сергек).",
        "Big Data жинау және өңдеу алгоритмдерін жасау.",
        "Бірегей LSTM нейрондық желісін оқыту және сынау.",
        "Мобильді платформада қосымшаны интеграциялау және AI моделін енгізу."
    ])

    add_content_slide("4. Зерттеу нысаны мен пәні", [
        "Нысаны: Қалалық көлік инфрақұрылымы мен ағындары.",
        "Пәні: Трафикті талдау және болжау алгоритмдері (ML/AI).",
        "Қолданылатын деректер: IoT сенсорлары, GPS тректері.",
        "Зерттеу аймағы: Астана қаласы (Есіл ауданы мысалында)."
    ])

    add_content_slide("5. Технологиялар Стегі", [
        "AI/ML Ядросы: PyTorch, Scikit-learn, авторлық LSTM алгоритмі.",
        "Деректер базасы: PostgreSQL, PostGIS (кеңістіктік индекстеу).",
        "Mobile: Flutter (Dart), Google Maps SDK.",
        "Серверлік инфрақұрылым: Асинхронды Python, Docker, Render."
    ])

    add_content_slide("6. Аналогтарды Талдау", [
        "2GIS / Yandex: Трафикті көрсетеді, бірақ AI-болжамы шектеулі, жабық архитектура.",
        "Сергек (Қазақстан): Жол ережесін бақылауға арналған, бірақ көлік ағынын алдын-ала болжамайды.",
        "AI Traffic: Ашық архитектура, терең нейрондық болжау, антистресс режим.",
        "Біздің артықшылық: Қазақстандық инфрақұрылымға 100% бейімделген бірегей AI алгоритмі."
    ])

    add_content_slide("7. Жүйе Архитектурасы", [
        "Көп деңгейлі архитектура: Clients -> Gateway -> AI Workers -> DB.",
        "WebSocket арқылы нақты уақыттағы деректер ағыны.",
        "Асинхронды өңдеу (Celery/Asyncio).",
        "Кеңістіктік индекстеу (R-tree/PostGIS) арқылы жылдам іздеу."
    ])

    add_content_slide("8. LSTM Алгоритмі және Қорытынды", [
        "Өзіміз оқытқан LSTM алгоритмі тарихи трафик пен ауа-райын талдап, MAE < 5% дәлдікпен болжайды.",
        "Жобаның негізгі құндылығы – дайын API-лерге емес, өзіміздің AI моделімізге негізделуінде.",
        "Дипломдық жұмыс мақсаты толық орындалып, қазақстандық нақты деректермен жұмыс істейтін AI-прототип жасалды.",
        "Жүйе қаланың көлік ағындарын оңтайландыруға және экологиялық жағдайын жақсартуға толық дайын."
    ])

    output_path = "AI_Traffic_Presentation_Final.pptx"
    prs.save(output_path)
    print(f"Created {output_path}")

if __name__ == "__main__":
    create_premium_pptx()
