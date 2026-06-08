import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333) # 16:9 ratio
prs.slide_height = Inches(7.5)

# Define corporate/academic colors
COLOR_PRIMARY = RGBColor(15, 32, 67)      # Deep Navy
COLOR_SECONDARY = RGBColor(0, 114, 198)   # Academic Blue
COLOR_ACCENT = RGBColor(230, 57, 70)      # Highlight Red
COLOR_TEXT = RGBColor(51, 51, 51)         # Dark Gray
COLOR_BG = RGBColor(250, 250, 250)        # Off-white

def add_footer(slide):
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(4), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "PhD Defense | L.N. Gumilyov ENU"
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(128, 128, 128)

def format_title(title_shape, text):
    title_shape.text = text
    p = title_shape.text_frame.paragraphs[0]
    p.font.name = 'Segoe UI'
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    p.alignment = PP_ALIGN.LEFT

def format_body(shape, text, size=18):
    shape.text = text
    for p in shape.text_frame.paragraphs:
        p.font.name = 'Segoe UI'
        p.font.size = Pt(size)
        p.font.color.rgb = COLOR_TEXT
        
def format_bullet(shape, texts, size=18):
    tf = shape.text_frame
    tf.clear()
    for text in texts:
        p = tf.add_paragraph()
        p.text = text
        p.font.name = 'Segoe UI'
        p.font.size = Pt(size)
        p.font.color.rgb = COLOR_TEXT
        p.level = 0

# Slide 1: Title
slide_layout = prs.slide_layouts[6] # blank
slide = prs.slides.add_slide(slide_layout)
bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
bg.fill.solid()
bg.fill.fore_color.rgb = COLOR_PRIMARY
bg.line.fill.background()

title = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11.333), Inches(2))
tf = title.text_frame
p = tf.paragraphs[0]
p.text = "Қалалық ортадағы көлік ағындарын бақылау мен болжауға арналған интеллектуалды жүйені архитектуралық жобалау"
p.font.name = 'Segoe UI'
p.font.size = Pt(40)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)
p.alignment = PP_ALIGN.CENTER

subtitle = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(11.333), Inches(1))
p2 = subtitle.text_frame.paragraphs[0]
p2.text = "PhD Dissertation Defense / Докторлық Диссертация Қорғау"
p2.font.name = 'Segoe UI'
p2.font.size = Pt(24)
p2.font.color.rgb = RGBColor(200, 200, 200)
p2.alignment = PP_ALIGN.CENTER

author = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(11.333), Inches(1.5))
p3 = author.text_frame.paragraphs[0]
p3.text = "Спикер: Сулейменов Алишер\nҒылыми жетекшісі: Кусаинова Айнұр\nЕҰУ им. Л.Н. Гумилева | Факультет Информационных Технологий"
p3.font.name = 'Segoe UI'
p3.font.size = Pt(18)
p3.font.color.rgb = RGBColor(255, 255, 255)
p3.alignment = PP_ALIGN.CENTER

# Slide 2: Relevance
slide = prs.slides.add_slide(slide_layout)
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(1))
format_title(title_box, "1. Зерттеудің Өзектілігі (Relevance & Motivation)")

content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(6), Inches(5))
bullets = [
    "• Қазақстанның ірі қалаларында көлік кептелісінің критикалық деңгейі (TomTom: 117 және 197 орындар).",
    "• Қолданыстағы жүйелердің (мысалы, 'Сергек') шектеулері:",
    "   - Тек айыппұл салуға және пост-фактум талдауға бағытталған.",
    "   - Кептелісті алдын ала динамикалық болжау мүмкіндігі жоқ.",
    "• Архитектуралық олқылықтар:",
    "   - Нақты уақыттағы 'Цифрлық Егіз' (Digital Twin) тұжырымдамасының жоқтығы.",
    "   - Ауа-райы мен инфрақұрылым факторларын кешенді есептемеу."
]
format_bullet(content_box, bullets, size=20)

if os.path.exists("traffic_trends.png"):
    slide.shapes.add_picture("traffic_trends.png", Inches(7), Inches(1.8), width=Inches(5.5))
add_footer(slide)

# Slide 3: Goals and Objectives
slide = prs.slides.add_slide(slide_layout)
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(1))
format_title(title_box, "2. Зерттеу Мақсаты мен Міндеттері")

content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(1.5))
tf = content_box.text_frame
p = tf.paragraphs[0]
p.text = "Мақсаты: LSTM нейрондық желісі мен Digital Twin тұжырымдамасы негізінде қалалық трафикті нақты уақыт режимінде бақылауға және болжауға арналған кешенді AI-жүйесін әзірлеу."
p.font.name = 'Segoe UI'
p.font.size = Pt(22)
p.font.bold = True
p.font.color.rgb = COLOR_SECONDARY

content_box2 = slide.shapes.add_textbox(Inches(0.5), Inches(3.0), Inches(12), Inches(4))
bullets = [
    "Ғылыми және Инженериялық Міндеттер:",
    "1. Жүйелік Архитектура: Микросервистік backend (FastAPI, PostGIS) және интеграциялық API шлюздерін жобалау.",
    "2. Машиналық Оқыту: Уақыттық қатарларды талдау үшін LSTM (PyTorch) моделін бейімдеу және оңтайландыру.",
    "3. Цифрлық Егіз: Кептелістің 'Digital Twin' моделін жасау және маршруттау алгоритмдерін (Dijkstra/A* модификациясы) іске асыру.",
    "4. Cross-Platform Клиент: Мобильді қосымша (Flutter) және Web Dashboard (Vue.js) арқылы UI/UX жобалау."
]
format_bullet(content_box2, bullets, size=20)
add_footer(slide)

# Slide 4: System Architecture
slide = prs.slides.add_slide(slide_layout)
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(1))
format_title(title_box, "3. Жүйелік Архитектура (System Architecture)")

content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(5), Inches(5))
bullets = [
    "• Микросервистік тәсіл:",
    "   - Масштабталатын API (FastAPI).",
    "   - Кеңістіктік деректер базасы (PostGIS).",
    "• Деректер ағыны:",
    "   - IoT сенсорлар мен GPS тректерден.",
    "   - Нақты уақыттағы өңдеу (Real-time).",
    "• Client-Server өзара әрекеттесуі (REST/WebSockets)."
]
format_bullet(content_box, bullets, size=18)

if os.path.exists("diag_architecture.png"):
    slide.shapes.add_picture("diag_architecture.png", Inches(5.5), Inches(1.5), width=Inches(7.5))
elif os.path.exists("component_diagram.png"):
    slide.shapes.add_picture("component_diagram.png", Inches(5.5), Inches(1.5), width=Inches(7.5))
add_footer(slide)

# Slide 5: ML Models & Algorithms
slide = prs.slides.add_slide(slide_layout)
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(1))
format_title(title_box, "4. Болжау Моделі: LSTM Нейрондық Желісі")

content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(6), Inches(5))
bullets = [
    "• Long Short-Term Memory (LSTM):",
    "   - Ұзақ мерзімді тәуелділіктерді сақтау мүмкіндігі.",
    "   - Кіріс: Соңғы 12 уақыттық нүкте (lookback).",
    "   - Шығыс: 30-60 минутқа трафик индексі.",
    "• Аномалияларды анықтау (Z-Score):",
    "   - Кептелістерді ерте ескерту.",
    "• Экспоненциалды жылжымалы орташа (EMA):",
    "   - α=0.4 коэффициентімен трендті тегістеу."
]
format_bullet(content_box, bullets, size=18)

if os.path.exists("lstm_architecture.png"):
    slide.shapes.add_picture("lstm_architecture.png", Inches(6.5), Inches(1.5), width=Inches(6.5))
elif os.path.exists("diag_lstm.png"):
    slide.shapes.add_picture("diag_lstm.png", Inches(6.5), Inches(1.5), width=Inches(6.5))
add_footer(slide)

# Slide 6: Results
slide = prs.slides.add_slide(slide_layout)
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(1))
format_title(title_box, "5. Модель Нәтижелері және Эксперимент (Results)")

content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(5.5), Inches(5))
bullets = [
    "• Метрикалар (Metrics):",
    "   - MAE (Mean Absolute Error) ~ 8% қателік.",
    "   - RMSE: Жалған аномалияларға төзімділік.",
    "• LSTM дәлдігі: 87.4%",
    "• Базалық модельдермен (Linear Regression) салыстырғанда дәлдіктің 34%-ға артуы.",
    "• Қосымша факторлар:",
    "   - Ауа-райы әсері математикалық түрде ескерілген."
]
format_bullet(content_box, bullets, size=18)

if os.path.exists("mae_rmse_chart.png"):
    slide.shapes.add_picture("mae_rmse_chart.png", Inches(6.5), Inches(1.5), width=Inches(6.5))
elif os.path.exists("model_comparison.png"):
    slide.shapes.add_picture("model_comparison.png", Inches(6.5), Inches(1.5), width=Inches(6.5))
add_footer(slide)

# Slide 7: Digital Twin
slide = prs.slides.add_slide(slide_layout)
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(1))
format_title(title_box, "6. Инновациялық Шешім: Digital Twin Симуляторы")

content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(5.5), Inches(5))
bullets = [
    "• What-If Симуляция (simulate_closure API):",
    "   - Жол жөндеу жұмыстары немесе апат (ДТП) жағдайында трафиктің қайта бөлінуін болжау.",
    "• Multimodal маршруттау:",
    "   - Автокөлік + Самокат/Жаяу жүру гибридті ұсыныстары.",
    "• Инклюзивті орта:",
    "   - Мүмкіндігі шектеулі жандарға арналған кедергісіз маршруттар (баспалдақтарды айналып өту)."
]
format_bullet(content_box, bullets, size=18)

if os.path.exists("diag_twin.png"):
    slide.shapes.add_picture("diag_twin.png", Inches(6.5), Inches(1.5), width=Inches(6.5))
elif os.path.exists("road_graph.png"):
    slide.shapes.add_picture("road_graph.png", Inches(6.5), Inches(1.5), width=Inches(6.5))
add_footer(slide)

# Slide 8: Conclusion
slide = prs.slides.add_slide(slide_layout)
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(1))
format_title(title_box, "7. Қорытынды және Практикалық Маңызы")

content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.333), Inches(5))
bullets = [
    "Ғылыми Жаңалығы:",
    "• Қазақстанның инфрақұрылымы мен климатына бейімделген LSTM-негізіндегі алғашқы трафик 'Цифрлық Егізі' құрылды.",
    "",
    "Экономикалық және Экологиялық Тиімділік:",
    "• Кептелістегі уақытты 20%-ға дейін қысқарту потенциалы.",
    "• CO2 шығарындыларын азайту арқылы қала экологиясын жақсарту.",
    "",
    "Әлеуметтік Маңызы:",
    "• Инклюзивті бағыттау және төтенше жағдайларда жедел әрекет етуді қамтамасыз ету."
]
format_bullet(content_box, bullets, size=22)
add_footer(slide)

# Slide 9: Q&A
slide = prs.slides.add_slide(slide_layout)
bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
bg.fill.solid()
bg.fill.fore_color.rgb = COLOR_PRIMARY
bg.line.fill.background()

title = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(11.333), Inches(2))
tf = title.text_frame
p = tf.paragraphs[0]
p.text = "Назарларыңызға рақмет!\nСұрақтарыңызға жауап беруге дайынмын."
p.font.name = 'Segoe UI'
p.font.size = Pt(40)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)
p.alignment = PP_ALIGN.CENTER

prs.save("AI_Traffic_Scopus_PhD_Defense.pptx")
print("Presentation saved successfully.")
