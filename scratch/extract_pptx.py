import sys
import codecs
try:
    import pptx
except ImportError:
    import subprocess
    subprocess.check_call([sys.executable, '-m', 'pip', 'install', 'python-pptx'])
    import pptx

with codecs.open('extracted_pptx_text.txt', 'w', encoding='utf-8') as f:
    prs = pptx.Presentation('AI_Traffic_Presentation_PhD_Edition.pptx')
    for i, slide in enumerate(prs.slides):
        f.write(f'--- Slide {i+1} ---\n')
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                f.write(shape.text + '\n')
