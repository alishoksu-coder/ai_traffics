from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def setup_title(slide, title_text):
    title = slide.shapes.title
    title.text = title_text
    title.text_frame.paragraphs[0].font.size = Pt(28)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

def create_dataset_presentation():
    prs = Presentation()
    
    # ----------------------------------------------------
    # Slide 10: Dataset логикасы
    # ----------------------------------------------------
    slide_10 = prs.slides.add_slide(prs.slide_layouts[5]) # blank with title
    setup_title(slide_10, "10-слайд. Dataset логикасы")
    
    # Left Side: Pipeline (Vertical)
    pipeline_steps = [
        "Historical / Simulated Source",
        "Validation",
        "Cleaning & Normalization",
        "Database: traffic_values",
        "ML input / Forecasting"
    ]
    
    left_x = Inches(0.5)
    start_y = Inches(1.5)
    box_width = Inches(3.5)
    box_height = Inches(0.6)
    spacing = Inches(0.9)
    
    for i, step in enumerate(pipeline_steps):
        # Add Box
        y = start_y + i * spacing
        shape = slide_10.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left_x, y, box_width, box_height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(79, 129, 189)
        shape.line.color.rgb = RGBColor(56, 93, 138)
        
        tf = shape.text_frame
        tf.text = step
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.bold = True
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(255, 255, 255)
        
        # Add Arrow (except for the last one)
        if i < len(pipeline_steps) - 1:
            arrow = slide_10.shapes.add_shape(
                MSO_SHAPE.DOWN_ARROW, 
                left_x + box_width/2 - Inches(0.15), 
                y + box_height + Inches(0.05), 
                Inches(0.3), 
                Inches(0.2)
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = RGBColor(150, 150, 150)
            arrow.line.color.rgb = RGBColor(150, 150, 150)

    # Right Side: Explanation Box
    right_x = Inches(4.5)
    explanation_width = Inches(5.0)
    
    exp_shape = slide_10.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, right_x, start_y, explanation_width, Inches(4.5)
    )
    exp_shape.fill.solid()
    exp_shape.fill.fore_color.rgb = RGBColor(240, 240, 240) # Light gray
    exp_shape.line.color.rgb = RGBColor(200, 200, 200)
    
    tf_exp = exp_shape.text_frame
    tf_exp.word_wrap = True
    
    p1 = tf_exp.paragraphs[0]
    p1.text = "Неге бұл тәсіл таңдалды?"
    p1.font.bold = True
    p1.font.size = Pt(16)
    p1.font.color.rgb = RGBColor(0, 51, 102)
    
    p2 = tf_exp.add_paragraph()
    p2.text = "Нақты қалалық ITS API қолжетімділігі шектеулі болғандықтан, жобада тестілеу және демонстрация үшін historical/simulated traffic dataset қолданылды. Бұл прототиптің шектеуі, бірақ архитектура real API интеграциясына дайын."
    p2.font.size = Pt(14)
    p2.font.color.rgb = RGBColor(50, 50, 50)
    p2.space_after = Pt(14)
    
    p3 = tf_exp.add_paragraph()
    p3.text = "Dataset мақсаты:"
    p3.font.bold = True
    p3.font.size = Pt(14)
    p3.font.color.rgb = RGBColor(0, 51, 102)
    
    p4 = tf_exp.add_paragraph()
    p4.text = "Жүйені тестілеу, ML-модельдерді салыстыру және dashboard көрсетілімін қамтамасыз ету."
    p4.font.size = Pt(14)
    p4.font.color.rgb = RGBColor(50, 50, 50)
    p4.space_after = Pt(20)
    
    p5 = tf_exp.add_paragraph()
    p5.text = "⚠️ Ескерту:"
    p5.font.bold = True
    p5.font.size = Pt(14)
    p5.font.color.rgb = RGBColor(192, 0, 0)
    
    p6 = tf_exp.add_paragraph()
    p6.text = "Бұл dataset прототиптік ортаға арналған. Келесі кезеңде real ITS / GPS / sensor data интеграциясы жоспарланады."
    p6.font.size = Pt(14)
    p6.font.color.rgb = RGBColor(50, 50, 50)


    # ----------------------------------------------------
    # Slide 11: Dataset белгілері
    # ----------------------------------------------------
    slide_11 = prs.slides.add_slide(prs.slide_layouts[5])
    setup_title(slide_11, "11-слайд. Dataset белгілері және олардың модельдегі рөлі")
    
    # Left: Table
    table_data = [
        ["Feature", "Рөлі"],
        ["timestamp / ts", "Уақыттық қатар реті"],
        ["location_id", "Жол нүктесі / учаске"],
        ["value / avg_speed_kmh", "Жүктеме немесе жылдамдық"],
        ["weather_factor", "Ауа райы әсері"],
        ["vehicle_count", "Көлік саны"],
        ["is_weekend", "Демалыс күн паттерні"],
        ["is_peak_hour", "Қарбалас уақыт"],
        ["accident_occurred / user_event", "Оқиға әсері"],
        ["congestion_level", "Target / болжанатын мән"]
    ]
    
    rows = len(table_data)
    cols = 2
    left = Inches(0.5)
    top = Inches(1.3)
    width = Inches(5.5)
    height = Inches(4.5)
    
    table_shape = slide_11.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    
    # Set column widths
    table.columns[0].width = Inches(2.5)
    table.columns[1].width = Inches(3.0)
    
    for row_idx, row_data in enumerate(table_data):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = cell_text
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(12)
            if row_idx == 0:
                p.font.bold = True
                p.font.color.rgb = RGBColor(255, 255, 255)
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(79, 129, 189)
            else:
                p.font.color.rgb = RGBColor(0, 0, 0)
    
    # Bottom Note
    note_top = top + height + Inches(0.2)
    note_box = slide_11.shapes.add_textbox(left, note_top, Inches(9.0), Inches(0.8))
    tf_note = note_box.text_frame
    tf_note.word_wrap = True
    p_note = tf_note.paragraphs[0]
    p_note.text = "💡 Маңызды: Нақты кодта негізгі өрістер location_id, ts, value, weather_factor түрінде сақталады. Қалған белгілер feature engineering кезеңінде қалыптастырылады."
    p_note.font.size = Pt(12)
    p_note.font.italic = True
    p_note.font.color.rgb = RGBColor(0, 102, 51)
    
    # Right: Feature Groups Sidebar
    right_x = Inches(6.5)
    card_width = Inches(3.0)
    card_height = Inches(1.1)
    card_spacing = Inches(1.3)
    
    groups = [
        {"title": "Уақыттық белгілер", "items": "timestamp, is_weekend, is_peak_hour", "color": RGBColor(220, 230, 242)},
        {"title": "Кеңістік белгілері", "items": "location_id, intersection_id", "color": RGBColor(228, 223, 236)},
        {"title": "Жағдайлық белгілер", "items": "weather, accident, vehicle_count", "color": RGBColor(235, 241, 222)}
    ]
    
    for i, grp in enumerate(groups):
        y = top + i * card_spacing
        shape = slide_11.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, right_x, y, card_width, card_height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = grp["color"]
        shape.line.color.rgb = RGBColor(150, 150, 150)
        
        tf = shape.text_frame
        tf.word_wrap = True
        
        p_title = tf.paragraphs[0]
        p_title.text = grp["title"]
        p_title.font.bold = True
        p_title.font.size = Pt(14)
        p_title.font.color.rgb = RGBColor(0, 51, 102)
        
        p_items = tf.add_paragraph()
        p_items.text = grp["items"]
        p_items.font.size = Pt(12)
        p_items.font.color.rgb = RGBColor(50, 50, 50)

    # ----------------------------------------------------
    # Slide 12: Data Flow
    # ----------------------------------------------------
    slide_12 = prs.slides.add_slide(prs.slide_layouts[5])
    setup_title(slide_12, "12-слайд. Жүйе ішіндегі Data Flow")
    
    flow_steps = [
        "Source\n(Simulator)",
        "Database\n(traffic_values)",
        "FastAPI\nendpoints",
        "ML\n(Forecasting)",
        "model_\nmetrics",
        "Dashboard\n(Digital Twin)"
    ]
    
    box_w = Inches(1.2)
    box_h = Inches(1.0)
    start_x = Inches(0.3)
    y_center = Inches(3.0)
    x_spacing = Inches(1.55)
    
    for i, step in enumerate(flow_steps):
        x = start_x + i * x_spacing
        
        shape = slide_12.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, y_center, box_w, box_h
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(31, 73, 125)
        shape.line.color.rgb = RGBColor(0, 0, 0)
        
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = step
        p.alignment = PP_ALIGN.CENTER
        p.font.bold = True
        p.font.size = Pt(11)
        p.font.color.rgb = RGBColor(255, 255, 255)
        
        if i < len(flow_steps) - 1:
            arrow = slide_12.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW, 
                x + box_w + Inches(0.05), 
                y_center + box_h/2 - Inches(0.1), 
                Inches(0.25), 
                Inches(0.2)
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = RGBColor(100, 100, 100)
            arrow.line.color.rgb = RGBColor(100, 100, 100)
            
    # Bottom Conclusion Box
    conc_box = slide_12.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(5.5), Inches(8.0), Inches(1.0)
    )
    conc_box.fill.solid()
    conc_box.fill.fore_color.rgb = RGBColor(240, 240, 240)
    conc_box.line.color.rgb = RGBColor(200, 200, 200)
    
    tf_conc = conc_box.text_frame
    tf_conc.word_wrap = True
    p_conc = tf_conc.paragraphs[0]
    p_conc.text = "Қорытынды:"
    p_conc.font.bold = True
    p_conc.font.size = Pt(14)
    p_conc.font.color.rgb = RGBColor(0, 51, 102)
    
    p_conc_desc = tf_conc.add_paragraph()
    p_conc_desc.text = "Дерек simulation/source қабатынан басталып, database, API, ML және dashboard модульдері арқылы толық өңделетін жүйелік pipeline құрайды."
    p_conc_desc.font.size = Pt(14)
    p_conc_desc.font.color.rgb = RGBColor(50, 50, 50)
    
    prs.save("ai_traffic_dataset_slides.pptx")
    print("Done generating ai_traffic_dataset_slides.pptx")

if __name__ == "__main__":
    create_dataset_presentation()
