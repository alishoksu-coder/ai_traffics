import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR

def create_phd_presentation():
    prs = Presentation()
    # Widescreen 16:9
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Academic Theme Colors
    BG_COLOR = RGBColor(250, 250, 250)      # Clean white background
    PRIMARY_TEXT = RGBColor(20, 30, 40)     # Dark slate
    SECONDARY_TEXT = RGBColor(80, 90, 100)  # Muted grey
    ACCENT_BLUE = RGBColor(0, 85, 164)      # Academic Blue
    ACCENT_RED = RGBColor(200, 30, 30)      # Alert Red
    BORDER_COLOR = RGBColor(220, 220, 220)

    # Helper to apply base styling
    def add_slide_with_title(title_text):
        slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank layout
        
        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = BG_COLOR
        
        # Header line
        header_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, Inches(0.8), prs.slide_width, Inches(0.05)
        )
        header_bar.fill.solid()
        header_bar.fill.fore_color.rgb = ACCENT_BLUE
        header_bar.line.fill.background()

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.6))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.bold = True
        p.font.size = Pt(28)
        p.font.color.rgb = ACCENT_BLUE
        
        return slide

    def add_textbox(slide, text, left, top, width, height, font_size=16, bold=False, color=PRIMARY_TEXT, align=PP_ALIGN.LEFT):
        txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.alignment = align
        return tf

    def add_bullet_points(tf, points, font_size=16):
        for idx, point in enumerate(points):
            if idx == 0:
                p = tf.paragraphs[0]
                p.text = point
            else:
                p = tf.add_paragraph()
                p.text = point
            p.font.size = Pt(font_size)
            p.font.color.rgb = SECONDARY_TEXT
            p.level = 0
            p.space_after = Pt(10)

    def set_notes(slide, text):
        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame
        text_frame.text = text

    # ==========================================
    # SLIDE 1: TITLE (PhD Defense Style)
    # ==========================================
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    slide1.background.fill.solid()
    slide1.background.fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    # Top border
    top_bar = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.2))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = ACCENT_BLUE
    top_bar.line.fill.background()

    add_textbox(slide1, "ЕҰУ им. Л.Н. Гумилева | Факультет Информационных Технологий", 
                1, 1, 11.33, 0.5, font_size=16, bold=True, color=SECONDARY_TEXT, align=PP_ALIGN.CENTER)
    
    add_textbox(slide1, "Қалалық ортадағы көлік ағындарын бақылау мен болжауға арналған интеллектуалды жүйені архитектуралық жобалау және әзірлеу", 
                1, 2.5, 11.33, 2, font_size=36, bold=True, color=PRIMARY_TEXT, align=PP_ALIGN.CENTER)

    add_textbox(slide1, "Специальность: «Вычислительная техника и программное обеспечение»", 
                1, 4.5, 11.33, 0.5, font_size=18, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)

    add_textbox(slide1, "Выполнил: Сулейменов Алишер\nНаучный руководитель: Кусаинова Айнур", 
                1, 5.5, 11.33, 1, font_size=16, bold=True, color=PRIMARY_TEXT, align=PP_ALIGN.CENTER)

    set_notes(slide1, "Қайырлы күн, құрметті комиссия мүшелері. Назарларыңызға «Қалалық ортадағы көлік ағындарын бақылау мен болжауға арналған интеллектуалды жүйе» тақырыбындағы зерттеу жұмысын ұсынамын.")

    # ==========================================
    # SLIDE 2: Актуальность & Аналоги (Merged per user request)
    # ==========================================
    slide2 = add_slide_with_title("1. Өзектілігі және қолданыстағы шешімдерді талдау (Relevance & Analogs)")
    
    # Left column: Text
    tf_relevance = add_textbox(slide2, "Мәселенің өзектілігі:", 0.5, 1.2, 5, 0.5, font_size=20, bold=True, color=PRIMARY_TEXT)
    add_bullet_points(tf_relevance, [
        "Қазақстанның ірі қалаларында (Астана, Алматы) көлік кептелісінің критикалық деңгейге жетуі (ТомТом бойынша 117-ші және 197-ші орындар).",
        "Қолданыстағы жүйелердің (Сергек) тек айыппұл мен пост-фактум анализге бағытталуы.",
        "Динамикалық болжау және «Цифрлық Егіз» (Digital Twin) архитектурасының жоқтығы."
    ], font_size=16)

    tf_analog = add_textbox(slide2, "Аналогтармен салыстыру:", 0.5, 3.5, 5, 0.5, font_size=20, bold=True, color=PRIMARY_TEXT)
    
    # Table for analogs
    rows = 4
    cols = 4
    left = Inches(0.5)
    top = Inches(4.2)
    width = Inches(6)
    height = Inches(2.5)
    
    table = slide2.shapes.add_table(rows, cols, left, top, width, height).table
    table.columns[0].width = Inches(1.5)
    table.columns[1].width = Inches(1.5)
    table.columns[2].width = Inches(1.5)
    table.columns[3].width = Inches(1.5)

    headers = ["Параметр", "2GIS", "Сергек", "AI Traffic (Біз)"]
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(14)

    data = [
        ("AI Болжау (60 мин)", "Шектеулі", "Жоқ", "Жоғары (LSTM)"),
        ("Аномалия детекторы", "Жоқ", "Жартылай", "Автоматты (Z-Score)"),
        ("Антистресс/Кедергісіз", "Жоқ", "Жоқ", "Толық қолдау")
    ]
    
    for row_idx, row_data in enumerate(data):
        for col_idx, text in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = text
            cell.text_frame.paragraphs[0].font.size = Pt(12)

    # Right column: Image (Traffic Map)
    add_textbox(slide2, "Астана қаласының ағымдағы кептеліс картасы (Мысал)", 7, 1.2, 5, 0.5, font_size=16, bold=True, align=PP_ALIGN.CENTER)
    if os.path.exists("road_graph.png"):
        slide2.shapes.add_picture("road_graph.png", Inches(7), Inches(1.8), width=Inches(5.5))
    elif os.path.exists("weather_impact.png"):
        slide2.shapes.add_picture("weather_impact.png", Inches(7), Inches(1.8), width=Inches(5.5))

    set_notes(slide2, "Қазіргі таңда трафик мәселесі өте өзекті. Экранның оң жағында көріп тұрғандарыңыздай, Астананың кептеліс картасында қызыл аймақтар көп. Біз нарықтағы 2GIS және Сергек сияқты аналогтармен нақты параметрлер бойынша салыстыру жүргіздік: олар болжауды және аномалияларды дәл біздің жүйедей терең талдамайды.")


    # ==========================================
    # SLIDE 3: Мақсат және Зерттеу міндеттері (Merged 2 and 4)
    # ==========================================
    slide3 = add_slide_with_title("2. Жұмыс мақсаты және зерттеу міндеттері")
    
    tf_goal = add_textbox(slide3, "Зерттеу мақсаты:", 0.5, 1.2, 12, 0.5, font_size=20, bold=True, color=PRIMARY_TEXT)
    add_textbox(slide3, "LSTM нейрондық желілік архитектурасы негізінде қалалық трафикті нақты уақыт режимінде мониторингтеуге және кептелісті алдын ала болжауға мүмкіндік беретін кешенді интеллектуалды AI-жүйесін әзірлеу.", 
                0.5, 1.7, 12, 1, font_size=18, color=SECONDARY_TEXT)

    tf_tasks = add_textbox(slide3, "Негізгі міндеттер (Architecture & Software Engineering):", 0.5, 3.2, 12, 0.5, font_size=20, bold=True, color=PRIMARY_TEXT)
    
    # Use smart shapes for tasks
    task_y = 4.0
    tasks = [
        "1. Backend (FastAPI, PostGIS) және интеграциялық API шлюздерін жобалау.",
        "2. Уақыттық қатарларды талдау үшін LSTM (PyTorch) моделін оқыту және оңтайландыру.",
        "3. Cross-platform Мобильді қосымша (Flutter) және Web Dashboard (Vue.js) жасау.",
        "4. Кептеліс «Цифрлық Егізін» (Digital Twin) және маршруттау алгоритмін (Dijkstra/A*) іске асыру."
    ]
    for i, t in enumerate(tasks):
        shape = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5 + (i%2)*6.2), Inches(task_y + (i//2)*1.2), Inches(6), Inches(0.9))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(240, 245, 250)
        shape.line.color.rgb = ACCENT_BLUE
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = t
        p.font.size = Pt(16)
        p.font.color.rgb = PRIMARY_TEXT
        p.alignment = PP_ALIGN.CENTER

    set_notes(slide3, "Зерттеуіміздің басты мақсаты – жай ғана карта жасау емес, LSTM алгоритмі арқылы трафикті болжайтын толыққанды архитектура құру. Ол үшін біз 4 негізгі міндетті қойдық: Серверлік бөлікті құру, AI моделін оқыту, Мобильді қосымша жасау және Цифрлық Егіз жүйесін іске асыру.")


    # ==========================================
    # SLIDE 4: Алгоритм функциялары және Архитектура (Side-by-side)
    # ==========================================
    slide4 = add_slide_with_title("3. Жүйе Архитектурасы және ML Алгоритмдерінің интеграциясы")
    
    add_textbox(slide4, "ML Болжау Алгоритмі (Функциялар мен Қызметтер)", 0.5, 1.2, 5.5, 0.5, font_size=18, bold=True, color=ACCENT_BLUE)
    
    tf_algo = add_textbox(slide4, "", 0.5, 1.8, 5.5, 5, font_size=14)
    add_bullet_points(tf_algo, [
        "predict_future(history, steps_ahead):",
        "  - Кіріс: Соңғы 12 уақыттық нүкте (lookback).",
        "  - Процесс: PyTorch тензорлары арқылы алға таралу (forward pass).",
        "  - Шығыс: Келесі 30-60 минутқа трафик индексі (0-100).",
        "",
        "predict_ema(series, alpha=0.4):",
        "  - Экспоненциалды жылжымалы орташа мән.",
        "  - Трафиктегі күрт өзгерістерге тез реакция береді.",
        "",
        "detect_anomaly(series):",
        "  - Соңғы 10 нүктедегі жылдамдықтың күрт төмендеуін талдайды (Z-Score негізінде).",
        "  - 'Critical' деңгейінде баламалы маршрут ұсынады."
    ], font_size=14)

    add_textbox(slide4, "Жүйелік Архитектура (Component Diagram)", 6.5, 1.2, 6, 0.5, font_size=18, bold=True, color=ACCENT_BLUE)
    
    if os.path.exists("diag_architecture.png"):
        slide4.shapes.add_picture("diag_architecture.png", Inches(6.5), Inches(1.8), width=Inches(6.3))
    elif os.path.exists("component_diagram.png"):
        slide4.shapes.add_picture("component_diagram.png", Inches(6.5), Inches(1.8), width=Inches(6.3))
    else:
        add_textbox(slide4, "[Архитектуралық схема: Flutter -> FastAPI -> PyTorch LSTM -> PostgreSQL]", 6.5, 3.5, 6, 1, font_size=16, align=PP_ALIGN.CENTER)

    set_notes(slide4, "Бұл слайдта жүйенің ішкі логикасы көрсетілген. Сол жақта басты функциялар: LSTM нейрондық желісі арқылы болжау, EMA арқылы күрт өзгерістерді бақылау және аномалияларды анықтау алгоритмдері. Оң жақта осы модульдердің микросервистік архитектурада қалай байланысқаны көрсетілген.")

    # ==========================================
    # SLIDE 5: Нәтижелер: Функциялар және Формулалар түсіндірмесі
    # ==========================================
    slide5 = add_slide_with_title("4. Модель нәтижелері және Математикалық негіздеме")
    
    add_textbox(slide5, "Метрикалардың математикалық мағынасы", 0.5, 1.2, 5.5, 0.5, font_size=18, bold=True, color=ACCENT_BLUE)
    
    tf_math = add_textbox(slide5, "", 0.5, 1.8, 5.5, 5, font_size=14)
    add_bullet_points(tf_math, [
        "MAE (Mean Absolute Error):",
        "  - Функция: mae_rmse()",
        "  - Қызметі: Болжам мен нақты деректің абсолютті ауытқуын есептеу.",
        "  - Түсіндірме: Егер MAE = 0.08 болса, модель трафик жүктемесін тек 8% қателікпен болжайды деген сөз.",
        "",
        "RMSE (Root Mean Square Error):",
        "  - Қызметі: Жалған және өте үлкен аномалияларды (шұғыл кептелістерді) қаттырақ жазалайтын метрика.",
        "",
        "Linear Regression Trend:",
        "  - Функция: predict_trend_lr()",
        "  - a = Σ(x-mx)(y-my) / Σ(x-mx)²",
        "  - Локальді бағытты (өсу/кему) анықтайды."
    ], font_size=14)

    add_textbox(slide5, "LSTM моделінің Эксперименттік Нәтижелері", 6.5, 1.2, 6, 0.5, font_size=18, bold=True, color=ACCENT_BLUE)
    
    # Table for metrics
    left = Inches(6.5)
    top = Inches(1.8)
    width = Inches(6)
    height = Inches(1.5)
    
    table2 = slide5.shapes.add_table(4, 4, left, top, width, height).table
    headers = ["Модель", "MAE", "RMSE", "Accuracy"]
    for i, h in enumerate(headers):
        table2.cell(0, i).text = h
        table2.cell(0, i).text_frame.paragraphs[0].font.bold = True
    
    data2 = [
        ("Linear Regression", "0.18", "0.24", "72.1%"),
        ("Random Forest", "0.11", "0.16", "81.5%"),
        ("LSTM (Ұсынылған)", "0.08", "0.12", "87.4%")
    ]
    
    for row_idx, row_data in enumerate(data2):
        for col_idx, text in enumerate(row_data):
            c = table2.cell(row_idx + 1, col_idx)
            c.text = text
            if row_idx == 2: # Highlight LSTM
                c.text_frame.paragraphs[0].font.bold = True
                c.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 128, 0)

    if os.path.exists("mae_rmse_chart.png"):
        slide5.shapes.add_picture("mae_rmse_chart.png", Inches(6.5), Inches(3.6), width=Inches(6.3))
    elif os.path.exists("model_comparison.png"):
        slide5.shapes.add_picture("model_comparison.png", Inches(6.5), Inches(3.6), width=Inches(6.3))

    set_notes(slide5, "Эксперимент нәтижелерін түсіндіретін болсақ: Сол жақта модель қателіктерін есептейтін математикалық аппарат пен функциялар көрсетілген. MAE метрикасы трафиктің абсолютті ауытқуын көрсетсе, RMSE үлкен аномалияларды анықтайды. Оң жақтағы кестеде біздің LSTM моделіміздің классикалық алгоритмдерге қарағанда әлдеқайда дәл (87.4%) екені дәлелденген.")

    # ==========================================
    # SLIDE 6: Цифрлық Егіз (Digital Twin) және Инновация
    # ==========================================
    slide6 = add_slide_with_title("5. Жүйенің инновациялық шешімдері: Digital Twin & Multimodal")
    
    tf_innov = add_textbox(slide6, "AI Traffic артықшылықтары:", 0.5, 1.2, 6, 0.5, font_size=18, bold=True)
    add_bullet_points(tf_innov, [
        "What-If Симуляторы (Digital Twin): Жол жөндеу немесе апат кезінде трафик ағынының қалай өзгеретінін алдын ала симуляциялау (simulate_closure API).",
        "Multimodal Analysis: Кептеліс критикалық деңгейге жеткенде, алгоритм автокөлікті қойып, самокатпен немесе жаяу жүруді ұсынады.",
        "Инклюзивті бағыттау: «Кедергісіз орта» режимі баспалдақтар мен кедергілерді айналып өтеді.",
        "Ауа-райы интеграциясы: Метео-деректер көлік ағынының жылдамдығына математикалық коэффициент (factor = 1.15 - 2.0) ретінде әсер етеді."
    ], font_size=15)

    if os.path.exists("diag_twin.png"):
        slide6.shapes.add_picture("diag_twin.png", Inches(7), Inches(1.5), width=Inches(5.5))
    elif os.path.exists("lstm_architecture.png"):
        slide6.shapes.add_picture("lstm_architecture.png", Inches(7), Inches(1.5), width=Inches(5.5))

    set_notes(slide6, "Жобаның басты ғылыми жаңалығы – Цифрлық Егіз технологиясы. Біз трафикті жай бақылап қана қоймай, жасанды түрде перекрытие жасап, ағынның қалай таралатынын симуляциялай аламыз. Сондай-ақ мультимодалды талдау кептеліс кезінде баламалы транспорт түрлерін ұсынады.")

    # ==========================================
    # SLIDE 7: Қорытынды және Практикалық маңызы
    # ==========================================
    slide7 = add_slide_with_title("6. Қорытынды және Практикалық Маңызы")
    
    tf_concl = add_textbox(slide7, "Жұмыс нәтижелері:", 0.5, 1.2, 12, 0.5, font_size=18, bold=True)
    add_bullet_points(tf_concl, [
        "PyTorch негізінде LSTM нейрондық желісі әзірленіп, дәлдігі 87.4%-ға жеткізілді.",
        "Микросервистік архитектура (FastAPI) және мобильді клиент (Flutter) толықтай іске қосылды.",
        "Қазақстанның климаттық және инфрақұрылымдық ерекшеліктерін ескеретін алгоритмдер құрылды."
    ], font_size=16)

    # Info boxes for impact
    shape1 = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(4), Inches(3.5), Inches(1.5))
    shape1.fill.solid()
    shape1.fill.fore_color.rgb = RGBColor(230, 240, 255)
    shape1.line.color.rgb = ACCENT_BLUE
    tf1 = shape1.text_frame
    tf1.text = "Экономикалық тиімділік\nКептелістегі уақытты 20%-ға қысқарту"
    tf1.paragraphs[0].font.size = Pt(16)
    tf1.paragraphs[0].font.bold = True
    tf1.paragraphs[0].font.color.rgb = ACCENT_BLUE
    tf1.paragraphs[0].alignment = PP_ALIGN.CENTER

    shape2 = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5), Inches(4), Inches(3.5), Inches(1.5))
    shape2.fill.solid()
    shape2.fill.fore_color.rgb = RGBColor(230, 255, 240)
    shape2.line.color.rgb = RGBColor(0, 150, 0)
    tf2 = shape2.text_frame
    tf2.text = "Экологиялық тиімділік\nCO2 шығарындыларын азайту"
    tf2.paragraphs[0].font.size = Pt(16)
    tf2.paragraphs[0].font.bold = True
    tf2.paragraphs[0].font.color.rgb = RGBColor(0, 100, 0)
    tf2.paragraphs[0].alignment = PP_ALIGN.CENTER

    shape3 = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9), Inches(4), Inches(3.5), Inches(1.5))
    shape3.fill.solid()
    shape3.fill.fore_color.rgb = RGBColor(255, 240, 240)
    shape3.line.color.rgb = ACCENT_RED
    tf3 = shape3.text_frame
    tf3.text = "Әлеуметтік маңыздылығы\nИнклюзивті орта және ДТП жедел реакциясы"
    tf3.paragraphs[0].font.size = Pt(16)
    tf3.paragraphs[0].font.bold = True
    tf3.paragraphs[0].font.color.rgb = ACCENT_RED
    tf3.paragraphs[0].alignment = PP_ALIGN.CENTER

    add_textbox(slide7, "Назарларыңызға рақмет!", 0.5, 6.2, 12, 1, font_size=28, bold=True, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)

    set_notes(slide7, "Қорытындылай келе, әзірленген жүйе тек теориялық зерттеу емес, нақты практикалық маңызға ие. Ол қала экономикасына, экологиясына және әлеуметтік инклюзиясына оң әсер етеді. Назарларыңызға рақмет, сұрақтарыңызға жауап беруге дайынмын.")

    prs.save("AI_Traffic_Presentation_PhD_Edition.pptx")
    print("Created AI_Traffic_Presentation_PhD_Edition.pptx successfully!")

if __name__ == "__main__":
    create_phd_presentation()
